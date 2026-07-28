from pathlib import Path

from drive_stripper import content


def test_read_text_plain_file(tmp_path: Path):
    path = tmp_path / "notes.txt"
    path.write_text("hello world")
    assert content.read_text(path) == "hello world"


def test_read_text_returns_none_for_image(tmp_path: Path):
    path = tmp_path / "photo.jpg"
    path.write_bytes(b"\xff\xd8\xff")  # not a valid jpeg, but suffix is enough for this check
    assert content.read_text(path) is None


def test_read_text_bim_file_treated_as_json_text(tmp_path: Path):
    path = tmp_path / "model.bim"
    path.write_text('{"name": "SalesModel", "owner": "jane@acme.com"}')
    assert content.read_text(path) == '{"name": "SalesModel", "owner": "jane@acme.com"}'
    assert content.is_text_extractable(path) is True


def test_write_docx_text_replaces_paragraph(tmp_path: Path):
    from docx import Document

    source = tmp_path / "doc.docx"
    dest = tmp_path / "doc.out.docx"
    doc = Document()
    doc.add_paragraph("contact jane@acme.com now")
    doc.save(source)

    content.write_docx_text(source, dest, lambda text: text.replace("jane@acme.com", "[REDACTED]"))

    result = Document(str(dest))
    assert result.paragraphs[0].text == "contact [REDACTED] now"


def test_write_xlsx_text_replaces_cell(tmp_path: Path):
    from openpyxl import Workbook, load_workbook

    source = tmp_path / "book.xlsx"
    dest = tmp_path / "book.out.xlsx"
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "jane@acme.com"
    wb.save(source)

    content.write_xlsx_text(source, dest, lambda text: text.replace("jane@acme.com", "[REDACTED]"))

    result = load_workbook(str(dest))
    assert result.active["A1"].value == "[REDACTED]"


def _make_pptx_fixture(path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(0, 0, 900000, 900000)
    textbox.text_frame.text = "contact jane@acme.com now"
    slide.shapes.add_table(1, 1, 0, 1000000, 2000000, 500000).table.cell(0, 0).text = (
        "cell jane@acme.com"
    )
    slide.notes_slide.notes_text_frame.text = "note jane@acme.com"
    prs.save(path)


def test_read_text_pptx_includes_shapes_tables_and_notes(tmp_path: Path):
    path = tmp_path / "deck.pptx"
    _make_pptx_fixture(path)

    text = content.read_text(path)
    assert text.count("jane@acme.com") == 3


def test_write_pptx_text_redacts_shapes_tables_and_notes(tmp_path: Path):
    from pptx import Presentation

    source = tmp_path / "deck.pptx"
    dest = tmp_path / "deck.out.pptx"
    _make_pptx_fixture(source)

    content.write_pptx_text(
        source, dest, lambda text: text.replace("jane@acme.com", "[REDACTED]")
    )

    prs = Presentation(str(dest))
    slide = prs.slides[0]
    texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
    table_text = slide.shapes[1].table.cell(0, 0).text
    notes_text = slide.notes_slide.notes_text_frame.text

    assert texts == ["contact [REDACTED] now"]
    assert table_text == "cell [REDACTED]"
    assert notes_text == "note [REDACTED]"
