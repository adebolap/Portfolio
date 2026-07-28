from pathlib import Path

import pytest

from drive_stripper import ocr, scaffold
from drive_stripper.pipeline import process_file


def test_strip_mode_on_text_file_redacts_and_reports_matches(tmp_path: Path):
    source = tmp_path / "notes.txt"
    source.write_text("email jane@acme.com about Project SkyNet")
    dest = tmp_path / "notes.sanitized.txt"

    result = process_file(source, dest, mode="strip", custom_terms=["Project SkyNet"])

    assert result.matches_found == 2
    text = dest.read_text()
    assert "jane@acme.com" not in text
    assert "Project SkyNet" not in text
    assert result.mapping_path is None


def test_scaffold_mode_on_text_file_is_reversible(tmp_path: Path):
    source = tmp_path / "notes.txt"
    original = "email jane@acme.com about the launch"
    source.write_text(original)
    dest = tmp_path / "notes.sanitized.txt"
    mapping_path = tmp_path / "notes.map.json"

    result = process_file(
        source, dest, mode="scaffold", mapping_destination=mapping_path
    )

    assert result.mapping_path == mapping_path
    scaffolded_text = dest.read_text()
    assert "jane@acme.com" not in scaffolded_text

    # simulate the round trip through a frontier model that echoes the text back unchanged
    model_response = scaffolded_text
    mapping = scaffold.load_mapping(mapping_path)
    restored = scaffold.restore(model_response, mapping)
    assert restored == original


def test_strip_mode_on_bim_file_redacts_json_text(tmp_path: Path):
    source = tmp_path / "model.bim"
    source.write_text('{"name": "SalesModel", "owner": "jane@acme.com"}')
    dest = tmp_path / "model.sanitized.bim"

    result = process_file(source, dest, mode="strip")

    assert result.matches_found == 1
    text = dest.read_text()
    assert "jane@acme.com" not in text
    assert "SalesModel" in text  # only the matched span is redacted, not the whole file


def test_strip_mode_on_pbids_file_redacts_connection_json(tmp_path: Path):
    source = tmp_path / "source.pbids"
    source.write_text('{"connections": [{"details": {"address": {"server": "10.0.4.22"}}}]}')
    dest = tmp_path / "source.sanitized.pbids"

    result = process_file(source, dest, mode="strip", categories=("ipv4",))

    assert result.matches_found == 1
    text = dest.read_text()
    assert "10.0.4.22" not in text
    assert "connections" in text


def test_strip_mode_on_pptx_redacts_slide_text(tmp_path: Path):
    from pptx import Presentation

    source = tmp_path / "deck.pptx"
    dest = tmp_path / "deck.sanitized.pptx"

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(0, 0, 900000, 900000).text_frame.text = "contact jane@acme.com"
    prs.save(source)

    result = process_file(source, dest, mode="strip", categories=("email",))

    assert result.matches_found == 1
    sanitized = Presentation(str(dest))
    shape_text = sanitized.slides[0].shapes[0].text_frame.text
    assert "jane@acme.com" not in shape_text


@pytest.mark.skipif(not ocr.is_available(), reason="tesseract OCR binary not installed")
def test_strip_mode_on_image_ocr_redacts_pixel_text(tmp_path: Path):
    from PIL import Image, ImageDraw, ImageFont

    source = tmp_path / "shot.png"
    dest = tmp_path / "shot.sanitized.png"
    font_path = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
    font = ImageFont.truetype(font_path, 28) if Path(font_path).exists() else ImageFont.load_default()
    img = Image.new("RGB", (700, 90), color="white")
    ImageDraw.Draw(img).text((10, 20), "Contact jane.doe@acme.com now", fill="black", font=font)
    img.save(source)

    result = process_file(source, dest, mode="strip", categories=("email",))

    assert result.image_regions_redacted >= 1
    assert result.ocr_warning is None
    assert "jane.doe@acme.com" not in ocr.extract_text(dest)


def test_scaffold_mode_on_image_raises(tmp_path: Path):
    source = tmp_path / "shot.png"
    from PIL import Image

    Image.new("RGB", (10, 10), color="white").save(source)
    dest = tmp_path / "shot.out.png"
    mapping_path = tmp_path / "shot.map.json"

    with pytest.raises(ValueError):
        process_file(source, dest, mode="scaffold", mapping_destination=mapping_path)


def test_metadata_only_mode_leaves_content_untouched(tmp_path: Path):
    source = tmp_path / "notes.txt"
    source.write_text("email jane@acme.com")
    dest = tmp_path / "notes.out.txt"

    result = process_file(source, dest, mode="metadata-only")

    assert result.matches_found == 0
    assert dest.read_text() == "email jane@acme.com"


def test_scaffold_mode_across_docx_paragraphs_keeps_tokens_unique(tmp_path: Path):
    from docx import Document

    source = tmp_path / "doc.docx"
    dest = tmp_path / "doc.sanitized.docx"
    mapping_path = tmp_path / "doc.map.json"

    doc = Document()
    doc.add_paragraph("reach jane@acme.com")
    doc.add_paragraph("or john@acme.com")
    doc.save(source)

    result = process_file(source, dest, mode="scaffold", mapping_destination=mapping_path)

    assert result.matches_found == 2
    mapping = scaffold.load_mapping(mapping_path)
    assert len(mapping) == 2  # tokens for both paragraphs are distinct, no collision

    sanitized = Document(str(dest))
    restored_paragraphs = [scaffold.restore(p.text, mapping) for p in sanitized.paragraphs]
    assert restored_paragraphs == ["reach jane@acme.com", "or john@acme.com"]
